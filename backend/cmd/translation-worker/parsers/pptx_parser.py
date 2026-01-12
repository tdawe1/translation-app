# parsers/pptx_parser.py
"""
PPTX parser using python-pptx for PowerPoint presentations.

Extracts text from slides with formatting and layout preservation.
"""

from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
import logging

logger = logging.getLogger(__name__)

# Try to import python-pptx
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt, Inches

    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False
    Presentation = None

# Import from plugins module
import sys
from pathlib import Path

worker_dir = Path(__file__).parent.parent
sys.path.insert(0, str(worker_dir))

from plugins import (
    ParserPlugin,
    ParsedDocument,
    Segment,
)


@dataclass
class SlideInfo:
    """Information about a PowerPoint slide."""

    number: int  # 1-based slide number
    layout_name: str = ""
    notes: str = ""  # Speaker notes
    has_title: bool = False
    shape_count: int = 0
    text_shape_count: int = 0


@dataclass
class TextFrameInfo:
    """Information about a text frame on a slide."""

    shape_id: int
    name: str = ""
    text: str = ""
    position: Tuple[float, float, float, float] = None  # (left, top, width, height)
    is_placeholder: bool = False
    has_text: bool = False

    font_size: float = 0.0
    font_name: str = ""
    bold: bool = False
    italic: bool = False
    underline: bool = False


@dataclass
class FormattingInfo:
    """Text formatting information."""

    font_name: str = ""
    font_size: float = 0.0
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color_rgb: Optional[Tuple[int, int, int]] = None
    alignment: str = ""  # left, center, right, justify


class PPTXParser:
    """PowerPoint PPTX parser using python-pptx.

    Features:
    - Extract text from slides with position context
    - Preserve formatting (bold, italic, font size)
    - Handle grouped and nested shapes
    - Extract speaker notes
    - Identify placeholder text
    - Support for master slides

    Note: python-pptx must be installed. Install with: pip install python-pptx
    """

    name = "pptx_parser"
    version = "1.0.0"
    dependencies = ["python-pptx"]

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """Initialize PPTX parser.

        Args:
            config: Optional configuration dict

        Raises:
            RuntimeError: If python-pptx is not installed
        """
        if not PYTHON_PPTX_AVAILABLE:
            raise RuntimeError(
                "python-pptx is not installed. Install with: pip install python-pptx"
            )

        self.config = config or {}
        self._extract_notes = self.config.get("extract_notes", False)
        self._extract_master = self.config.get("extract_master", False)
        self._merge_paragraphs = self.config.get("merge_paragraphs", True)
        self._include_empty_slides = self.config.get("include_empty_slides", False)

    def initialize(self, config: Dict[str, Any]) -> None:
        """Initialize plugin with configuration.

        Args:
            config: Configuration dict
        """
        self.config.update(config)
        self._extract_notes = self.config.get("extract_notes", False)
        self._extract_master = self.config.get("extract_master", False)
        self._merge_paragraphs = self.config.get("merge_paragraphs", True)
        self._include_empty_slides = self.config.get("include_empty_slides", False)
        logger.info(f"[{self.name}] Initialized with config: {self.config}")

    def shutdown(self) -> None:
        """Clean up resources."""
        logger.debug(f"[{self.name}] Shutdown")

    def supported_extensions(self) -> List[str]:
        """Return list of supported file extensions.

        Returns:
            List of supported extensions
        """
        return [".pptx", ".ppt"]

    def parse(self, file_path: str) -> ParsedDocument:
        """Parse PPTX and extract translatable text segments.

        Args:
            file_path: Path to PPTX file

        Returns:
            ParsedDocument with segments

        Raises:
            FileNotFoundError: If file doesn't exist
            Exception: If PPTX cannot be parsed
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {file_path}")

        prs = Presentation(file_path)

        try:
            segments = []
            metadata = {
                "format": "pptx",
                "slide_count": len(prs.slides),
                "title": self._get_presentation_title(prs),
            }

            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_info = self._extract_slide_info(slide, slide_idx)

                # Extract text frames
                text_frames = self._extract_text_frames(slide)

                # Check if slide has any translatable content
                has_content = any(tf.has_text for tf in text_frames)

                if not has_content and not self._include_empty_slides:
                    continue

                # Create segments from text frames
                slide_segments = self._create_slide_segments(slide_info, text_frames)
                segments.extend(slide_segments)

                # Add notes segment if requested
                if self._extract_notes and slide_info.notes:
                    segments.append(
                        Segment(
                            id=f"slide_{slide_idx}_notes",
                            text=slide_info.notes,
                            context={"type": "notes", "slide_number": slide_idx},
                            metadata={"is_notes": True},
                        )
                    )

            return ParsedDocument(
                segments=segments,
                metadata=metadata,
                format="pptx",
                source_path=file_path,
            )

        except Exception as e:
            logger.error(f"[{self.name}] Failed to parse PPTX: {e}")
            raise
        finally:
            # Presentation doesn't need explicit closing in python-pptx
            pass

    def _get_presentation_title(self, prs) -> str:
        """Extract presentation title from first slide or metadata.

        Args:
            prs: Presentation object

        Returns:
            Presentation title
        """
        if prs.slides:
            first_slide = prs.slides[0]
            for shape in first_slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text.strip()
                    if text and len(text) < 100:  # Likely title
                        return text

        # Fallback to filename
        return prs.core.title if hasattr(prs.core, "title") else ""

    def _extract_slide_info(self, slide, slide_idx: int) -> SlideInfo:
        """Extract information about a slide.

        Args:
            slide: python-pptx Slide object
            slide_idx: 1-based slide number

        Returns:
            SlideInfo with slide details
        """
        shape_count = len(slide.shapes)
        text_shape_count = sum(
            1 for s in slide.shapes if hasattr(s, "text_frame") and s.text_frame
        )

        # Check for title placeholder
        has_title = any(
            getattr(s, "is_placeholder", False)
            for s in slide.shapes
            if hasattr(s, "shape_type")
        )

        # Extract notes
        notes = ""
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text

        return SlideInfo(
            number=slide_idx,
            layout_name=getattr(slide, "name", ""),
            notes=notes,
            has_title=has_title,
            shape_count=shape_count,
            text_shape_count=text_shape_count,
        )

    def _extract_text_frames(self, slide) -> List[TextFrameInfo]:
        """Extract all text frames from a slide.

        Args:
            slide: python-pptx Slide object

        Returns:
            List of TextFrameInfo objects
        """
        text_frames = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Recursively extract from grouped shapes
                group_frames = self._extract_from_group(shape)
                text_frames.extend(group_frames)
                continue

            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue

            text_frame = shape.text_frame
            if not text_frame.text:
                continue

            # Get frame position
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            # Extract formatting from first run
            formatting = self._extract_formatting(text_frame)

            text_frames.append(
                TextFrameInfo(
                    shape_id=shape.shape_id,
                    name=shape.name,
                    text=text_frame.text,
                    position=(left, top, width, height),
                    is_placeholder=getattr(shape, "is_placeholder", False),
                    has_text=True,
                    font_size=formatting.font_size,
                    font_name=formatting.font_name,
                    bold=formatting.bold,
                    italic=formatting.italic,
                    underline=formatting.underline,
                )
            )

        return text_frames

    def _extract_from_group(self, group_shape) -> List[TextFrameInfo]:
        """Extract text frames from a grouped shape.

        Args:
            group_shape: Group shape object

        Returns:
            List of TextFrameInfo objects
        """
        text_frames = []

        for shape in group_shape.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                text_frames.extend(self._extract_from_group(shape))
            elif hasattr(shape, "text_frame") and shape.text_frame:
                if shape.text_frame.text:
                    formatting = self._extract_formatting(shape.text_frame)
                    text_frames.append(
                        TextFrameInfo(
                            shape_id=shape.shape_id,
                            name=shape.name,
                            text=shape.text_frame.text,
                            is_placeholder=getattr(shape, "is_placeholder", False),
                            has_text=True,
                            font_size=formatting.font_size,
                            font_name=formatting.font_name,
                            bold=formatting.bold,
                            italic=formatting.italic,
                            underline=formatting.underline,
                        )
                    )

        return text_frames

    def _extract_formatting(self, text_frame) -> FormattingInfo:
        """Extract formatting from text frame.

        Args:
            text_frame: python-pptx TextFrame object

        Returns:
            FormattingInfo with formatting details
        """
        info = FormattingInfo()

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font:
                    info.font_name = run.font.name
                    info.bold = run.font.bold
                    info.italic = run.font.italic
                    info.underline = run.font.underline

                    size_pt = run.font.size
                    if size_pt:
                        info.font_size = size_pt

                    # Extract color if available
                    if run.font.color and run.font.color.type == 1:  # RGB
                        info.color_rgb = (
                            run.font.color.rgb,
                            run.font.color.rgb,
                            run.font.color.rgb,
                        )

                break  # Use first run's formatting

        return info

    def _create_slide_segments(
        self, slide_info: SlideInfo, text_frames: List[TextFrameInfo]
    ) -> List[Segment]:
        """Create segments from slide text frames.

        Args:
            slide_info: Information about the slide
            text_frames: List of text frames on the slide

        Returns:
            List of Segment objects
        """
        segments = []

        if self._merge_paragraphs:
            # Merge all text on slide into one segment
            merged_text = "\n".join(tf.text for tf in text_frames if tf.text)

            if merged_text:
                segments.append(
                    Segment(
                        id=f"slide_{slide_info.number}",
                        text=merged_text,
                        context={
                            "type": "slide",
                            "slide_number": slide_info.number,
                            "layout_name": slide_info.layout_name,
                        },
                        metadata={
                            "shape_count": slide_info.shape_count,
                            "text_shape_count": slide_info.text_shape_count,
                            "has_title": slide_info.has_title,
                        },
                    )
                )
        else:
            # Create segment per text frame
            for idx, text_frame in enumerate(text_frames):
                segments.append(
                    Segment(
                        id=f"slide_{slide_info.number}_shape_{idx}",
                        text=text_frame.text,
                        context={
                            "type": "text_frame",
                            "slide_number": slide_info.number,
                            "shape_id": text_frame.shape_id,
                            "shape_name": text_frame.name,
                        },
                        metadata={
                            "position": text_frame.position,
                            "font_size": text_frame.font_size,
                            "font_name": text_frame.font_name,
                            "bold": text_frame.bold,
                            "italic": text_frame.italic,
                            "underline": text_frame.underline,
                            "is_placeholder": text_frame.is_placeholder,
                        },
                    )
                )

        return segments

    def render(
        self, doc: ParsedDocument, output_path: str, template_path: Optional[str] = None
    ) -> None:
        """Render translated PPTX with original layout.

        Args:
            doc: ParsedDocument with translated segments
            output_path: Where to save the output PPTX
            template_path: Optional original PPTX as template

        Raises:
            Exception: If rendering fails
        """
        if not PYTHON_PPTX_AVAILABLE:
            raise RuntimeError("python-pptx is not installed")

        # Load template or create new presentation
        if template_path:
            try:
                prs = Presentation(template_path)
            except Exception as e:
                logger.warning(f"[{self.name}] Failed to load template: {e}")
                prs = Presentation()
        else:
            prs = Presentation()

        try:
            # Map segments by slide number
            slide_map: Dict[int, List[Segment]] = {}
            for segment in doc.segments:
                slide_num = segment.context.get("slide_number", 1)
                if slide_num not in slide_map:
                    slide_map[slide_num] = []
                slide_map[slide_num].append(segment)

            # Process each slide
            for slide_num, segments in slide_map.items():
                # Adjust to 0-based index
                slide_idx = slide_num - 1

                # Ensure we have enough slides
                while len(prs.slides) <= slide_idx:
                    prs.slides.add_slide(prs.slide_layouts[0])

                slide = prs.slides[slide_idx]

                # Clear existing text (in production, would use more sophisticated approach)
                self._clear_slide_text(slide)

                # Insert translated text
                for segment in segments:
                    self._insert_segment_text(slide, segment)

            prs.save(output_path)
            logger.info(f"[{self.name}] Rendered PPTX to: {output_path}")

        except Exception as e:
            logger.error(f"[{self.name}] Failed to render PPTX: {e}")
            raise
        finally:
            # Presentation auto-saves, no explicit close needed
            pass

    def _clear_slide_text(self, slide) -> None:
        """Clear all text from a slide.

        Args:
            slide: python-pptx Slide object
        """
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                # Clear paragraph by paragraph to preserve structure
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.text = ""

    def _insert_segment_text(self, slide, segment: Segment) -> None:
        """Insert translated text into slide.

        Args:
            slide: python-pptx Slide object
            segment: Segment with translated text
        """
        seg_type = segment.context.get("type", "slide")

        if seg_type == "notes":
            # Insert into notes
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                slide.notes_slide.notes_text_frame.text = segment.text
            return

        # For regular slide text, try to find matching text box
        shape_name = segment.context.get("shape_name", "")
        position = segment.metadata.get("position")

        # Try to find existing text box by position
        target_shape = None
        if position:
            target_shape = self._find_shape_at_position(slide, position)

        # If no match, add new text box
        if target_shape is None:
            left = top = 100  # Default position
            if position:
                left, top = position[0], position[1]

            textbox = slide.shapes.add_textbox(
                left, top, width=Inches(6), height=Inches(1)
            )
            target_shape = textbox

        # Insert text
        if target_shape.has_text_frame:
            text_frame = target_shape.text_frame
            text_frame.clear()  # Clear existing

            # Apply formatting if available
            for paragraph in segment.text.split("\n"):
                p = text_frame.add_paragraph()
                run = (
                    p.add_run()
                )  # Create a run first, add_paragraph() creates empty paragraph
                run.text = paragraph

                # Apply formatting from metadata
                if segment.metadata.get("bold"):
                    run.font.bold = True
                if segment.metadata.get("italic"):
                    run.font.italic = True
                if segment.metadata.get("font_size"):
                    run.font.size = Pt(segment.metadata["font_size"])
                if segment.metadata.get("font_name"):
                    run.font.name = segment.metadata["font_name"]

    def _find_shape_at_position(self, slide, position) -> Optional:
        """Find a shape at or near the given position.

        Args:
            slide: python-pptx Slide object
            position: (left, top, width, height) tuple

        Returns:
            Matching shape or None
        """
        left, top, width, height = position
        tolerance = 10  # pixels

        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            if abs(shape.left - left) < tolerance and abs(shape.top - top) < tolerance:
                return shape

        return None

    def get_slide_count(self, file_path: str) -> int:
        """Get the number of slides in a presentation.

        Args:
            file_path: Path to PPTX file

        Returns:
            Number of slides
        """
        if not PYTHON_PPTX_AVAILABLE:
            raise RuntimeError("python-pptx is not installed")

        prs = Presentation(file_path)
        return len(prs.slides)

    def extract_slide_titles(self, file_path: str) -> List[str]:
        """Extract titles from all slides.

        Args:
            file_path: Path to PPTX file

        Returns:
            List of slide titles (empty string for slides without titles)
        """
        prs = Presentation(file_path)
        titles = []

        for slide in prs.slides:
            title = ""
            # Try to find title placeholder
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text and text[:30]:  # Assume short text is title
                        title = text
                        break
            titles.append(title)

        return titles


def create_pptx_parser(config: Optional[Dict[str, Any]] = None) -> PPTXParser:
    """Factory function to create a PPTX parser.

    Args:
        config: Optional configuration dict

    Returns:
        PPTXParser instance

    Raises:
        RuntimeError: If python-pptx is not installed
    """
    return PPTXParser(config)
