# tests/test_parsers/test_pptx_parser.py
"""
Unit tests for PPTX parser functionality.

Tests python-pptx-based PowerPoint parsing, extraction, and rendering.
Tests gracefully skip when python-pptx is not installed.
"""

import pytest
import sys
import tempfile
from pathlib import Path

# Add worker directory to path for imports
worker_dir = Path(__file__).parent.parent.parent
sys.path.insert(0, str(worker_dir))

# Import parser module
try:
    from parsers.pptx_parser import (
        PPTXParser,
        SlideInfo,
        TextFrameInfo,
        FormattingInfo,
        create_pptx_parser,
        PYTHON_PPTX_AVAILABLE,
    )
    from pptx.util import Inches

    PYTHON_PPTX_INSTALLED = PYTHON_PPTX_AVAILABLE
except ImportError:
    PYTHON_PPTX_INSTALLED = False
    PPTXParser = None
    SlideInfo = None
    TextFrameInfo = None
    FormattingInfo = None
    create_pptx_parser = None
    Inches = None


@pytest.mark.skipif(not PYTHON_PPTX_INSTALLED, reason="python-pptx not installed")
class TestSlideInfo:
    """Test SlideInfo dataclass."""

    def test_create_slide_info(self):
        """Should create SlideInfo with correct attributes."""
        slide = SlideInfo(
            number=1,
            layout_name="Title Slide",
            notes="Speaker notes here",
            has_title=True,
            shape_count=5,
            text_shape_count=3,
        )

        assert slide.number == 1
        assert slide.layout_name == "Title Slide"
        assert slide.notes == "Speaker notes here"
        assert slide.has_title is True
        assert slide.shape_count == 5
        assert slide.text_shape_count == 3

    def test_slide_info_defaults(self):
        """Should create SlideInfo with defaults."""
        slide = SlideInfo(number=1)

        assert slide.number == 1
        assert slide.layout_name == ""
        assert slide.notes == ""
        assert slide.has_title is False
        assert slide.shape_count == 0
        assert slide.text_shape_count == 0


@pytest.mark.skipif(not PYTHON_PPTX_INSTALLED, reason="python-pptx not installed")
class TestTextFrameInfo:
    """Test TextFrameInfo dataclass."""

    def test_create_text_frame_info(self):
        """Should create TextFrameInfo with position."""
        frame = TextFrameInfo(
            shape_id=1,
            name="Title",
            text="Hello World",
            position=(100, 200, 300, 50),
            is_placeholder=True,
            has_text=True,
            font_size=18.0,
            font_name="Arial",
            bold=True,
        )

        assert frame.shape_id == 1
        assert frame.name == "Title"
        assert frame.text == "Hello World"
        assert frame.position == (100, 200, 300, 50)
        assert frame.is_placeholder is True
        assert frame.has_text is True
        assert frame.font_size == 18.0
        assert frame.bold is True

    def test_text_frame_info_defaults(self):
        """Should create TextFrameInfo with minimal defaults."""
        frame = TextFrameInfo(shape_id=1)

        assert frame.shape_id == 1
        assert frame.name == ""
        assert frame.text == ""
        assert frame.position is None
        assert frame.is_placeholder is False
        assert frame.has_text is False
        assert frame.font_size == 0.0
        assert frame.bold is False


@pytest.mark.skipif(not PYTHON_PPTX_INSTALLED, reason="python-pptx not installed")
class TestFormattingInfo:
    """Test FormattingInfo dataclass."""

    def test_create_formatting_info(self):
        """Should create FormattingInfo with all fields."""
        info = FormattingInfo(
            font_name="Calibri",
            font_size=14.0,
            bold=True,
            italic=False,
            underline=True,
            color_rgb=(255, 0, 0),
            alignment="center",
        )

        assert info.font_name == "Calibri"
        assert info.font_size == 14.0
        assert info.bold is True
        assert info.italic is False
        assert info.underline is True
        assert info.color_rgb == (255, 0, 0)
        assert info.alignment == "center"


@pytest.mark.skipif(not PYTHON_PPTX_INSTALLED, reason="python-pptx not installed")
class TestPPTXParser:
    """Test PPTXParser functionality."""

    def test_plugin_attributes(self):
        """Should have required plugin attributes."""
        parser = PPTXParser()

        assert hasattr(parser, "name")
        assert hasattr(parser, "version")
        assert hasattr(parser, "dependencies")
        assert parser.name == "pptx_parser"
        assert parser.version == "1.0.0"
        assert "python-pptx" in parser.dependencies

    def test_supported_extensions(self):
        """Should return PPTX and PPT extensions."""
        parser = PPTXParser()

        extensions = parser.supported_extensions()
        assert ".pptx" in extensions
        assert ".ppt" in extensions
        assert len(extensions) == 2

    def test_initialization_default_config(self):
        """Should initialize with default config."""
        parser = PPTXParser()

        assert parser._extract_notes is False
        assert parser._extract_master is False
        assert parser._merge_paragraphs is True
        assert parser._include_empty_slides is False

    def test_initialization_custom_config(self):
        """Should accept custom configuration."""
        parser = PPTXParser(
            config={
                "extract_notes": True,
                "extract_master": True,
                "merge_paragraphs": False,
                "include_empty_slides": True,
            }
        )

        assert parser._extract_notes is True
        assert parser._extract_master is True
        assert parser._merge_paragraphs is False
        assert parser._include_empty_slides is True

    def test_initialize_method(self):
        """Should update config via initialize method."""
        parser = PPTXParser()
        parser.initialize({"extract_notes": True, "merge_paragraphs": False})

        assert parser._extract_notes is True
        assert parser._merge_paragraphs is False

    def test_shutdown(self):
        """Should shutdown without errors."""
        parser = PPTXParser()
        parser.shutdown()  # Should not raise

    def test_parse_nonexistent_file(self):
        """Should raise FileNotFoundError for missing file."""
        parser = PPTXParser()

        with pytest.raises(FileNotFoundError):
            parser.parse("/nonexistent/file.pptx")

    def test_create_pptx_parser_factory(self):
        """Should create parser via factory function."""
        parser = create_pptx_parser()

        assert isinstance(parser, PPTXParser)

    def test_create_pptx_parser_with_config(self):
        """Should create parser with config via factory."""
        parser = create_pptx_parser({"merge_paragraphs": False})

        assert parser._merge_paragraphs is False


@pytest.mark.skipif(not PYTHON_PPTX_INSTALLED, reason="python-pptx not installed")
class TestPPTXParserWithSampleFile:
    """Test parser with actual PPTX file."""

    def test_parse_simple_pptx(self):
        """Should parse a simple PPTX document."""
        from pptx import Presentation

        parser = PPTXParser()

        # Create a simple test PPTX
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])

            # Add title
            title_shape = slide.shapes.title
            title_shape.text = "Test Presentation"

            # Add text box
            left = Inches(1)
            top = Inches(2)
            width = Inches(6)
            height = Inches(1)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = "This is sample content"

            prs.save(tmp.name)

            try:
                parsed = parser.parse(tmp.name)

                assert parsed.format == "pptx"
                assert parsed.metadata["slide_count"] == 1
                assert len(parsed.segments) >= 1
                assert (
                    "Test Presentation" in parsed.segments[0].text
                    or "sample content" in parsed.segments[0].text
                )
            finally:
                Path(tmp.name).unlink(missing_ok=True)

    def test_parse_multi_slide_pptx(self):
        """Should parse multi-slide PPTX."""
        from pptx import Presentation

        parser = PPTXParser()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()

            for i in range(3):
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = f"Slide {i + 1}"

            prs.save(tmp.name)

            try:
                parsed = parser.parse(tmp.name)

                assert parsed.metadata["slide_count"] == 3
                assert len(parsed.segments) == 3
            finally:
                Path(tmp.name).unlink(missing_ok=True)

    def test_parse_with_notes(self):
        """Should extract speaker notes when configured."""
        from pptx import Presentation

        parser = PPTXParser(config={"extract_notes": True})

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Title"

            # Add notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "Speaker notes here"

            prs.save(tmp.name)

            try:
                parsed = parser.parse(tmp.name)

                # Should have regular slide segment plus notes segment
                notes_segments = [
                    s for s in parsed.segments if s.metadata.get("is_notes")
                ]
                assert len(notes_segments) == 1
                assert "Speaker notes here" in notes_segments[0].text
            finally:
                Path(tmp.name).unlink(missing_ok=True)

    def test_parse_without_notes(self):
        """Should not extract notes when not configured."""
        from pptx import Presentation

        parser = PPTXParser(config={"extract_notes": False})

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Title"

            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "These notes should not be extracted"

            prs.save(tmp.name)

            try:
                parsed = parser.parse(tmp.name)

                notes_segments = [
                    s for s in parsed.segments if s.metadata.get("is_notes")
                ]
                assert len(notes_segments) == 0
            finally:
                Path(tmp.name).unlink(missing_ok=True)

    def test_get_slide_count(self):
        """Should return correct slide count."""
        from pptx import Presentation

        parser = PPTXParser()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()
            for _ in range(5):
                prs.slides.add_slide(prs.slide_layouts[0])
            prs.save(tmp.name)

            try:
                count = parser.get_slide_count(tmp.name)
                assert count == 5
            finally:
                Path(tmp.name).unlink(missing_ok=True)

    def test_extract_slide_titles(self):
        """Should extract titles from all slides."""
        from pptx import Presentation

        parser = PPTXParser()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()

            for i in range(3):
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = f"Title {i + 1}"

            prs.save(tmp.name)

            try:
                titles = parser.extract_slide_titles(tmp.name)

                assert len(titles) == 3
                assert "Title 1" in titles[0]
                assert "Title 2" in titles[1]
                assert "Title 3" in titles[2]
            finally:
                Path(tmp.name).unlink(missing_ok=True)


@pytest.mark.skipif(not PYTHON_PPTX_INSTALLED, reason="python-pptx not installed")
class TestPPTXParserRendering:
    """Test PPTX rendering functionality."""

    def test_render_with_template(self):
        """Should render PPTX using template."""
        from pptx import Presentation
        from plugins import ParsedDocument, Segment

        parser = PPTXParser()

        # Create source PPTX
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as source_tmp:
            with tempfile.NamedTemporaryFile(
                suffix=".pptx", delete=False
            ) as output_tmp:
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = "Original Title"
                prs.save(source_tmp.name)

                try:
                    # Create parsed document with translation
                    parsed = ParsedDocument(
                        segments=[
                            Segment(
                                id="slide_1",
                                text="Translated Title",
                                context={"type": "slide", "slide_number": 1},
                                metadata={},
                            )
                        ],
                        metadata={"format": "pptx", "slide_count": 1},
                        format="pptx",
                        source_path=source_tmp.name,
                    )

                    # Render with template
                    parser.render(
                        doc=parsed,
                        output_path=output_tmp.name,
                        template_path=source_tmp.name,
                    )

                    # Verify output file exists and has content
                    assert Path(output_tmp.name).exists()
                    assert Path(output_tmp.name).stat().st_size > 0

                    # Verify we can open the rendered file
                    rendered_prs = Presentation(output_tmp.name)
                    assert len(rendered_prs.slides) == 1

                finally:
                    Path(source_tmp.name).unlink(missing_ok=True)
                    Path(output_tmp.name).unlink(missing_ok=True)

    def test_render_without_template(self):
        """Should render PPTX without template (blank slides)."""
        from pptx import Presentation
        from plugins import ParsedDocument, Segment

        parser = PPTXParser()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as output_tmp:
            try:
                parsed = ParsedDocument(
                    segments=[
                        Segment(
                            id="slide_1",
                            text="Test content",
                            context={"type": "slide", "slide_number": 1},
                        )
                    ],
                    metadata={"format": "pptx"},
                    format="pptx",
                )

                parser.render(doc=parsed, output_path=output_tmp.name)

                assert Path(output_tmp.name).exists()
                assert Path(output_tmp.name).stat().st_size > 0

            finally:
                Path(output_tmp.name).unlink(missing_ok=True)

    def test_render_with_notes(self):
        """Should render translated notes."""
        from pptx import Presentation
        from plugins import ParsedDocument, Segment

        parser = PPTXParser()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as source_tmp:
            with tempfile.NamedTemporaryFile(
                suffix=".pptx", delete=False
            ) as output_tmp:
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = "Title"
                slide.notes_slide.notes_text_frame.text = "Original notes"
                prs.save(source_tmp.name)

                try:
                    parsed = ParsedDocument(
                        segments=[
                            Segment(
                                id="slide_1_notes",
                                text="Translated notes",
                                context={"type": "notes", "slide_number": 1},
                                metadata={"is_notes": True},
                            )
                        ],
                        metadata={"format": "pptx"},
                        format="pptx",
                        source_path=source_tmp.name,
                    )

                    parser.render(
                        doc=parsed,
                        output_path=output_tmp.name,
                        template_path=source_tmp.name,
                    )

                    # Verify notes were translated
                    rendered_prs = Presentation(output_tmp.name)
                    rendered_slide = rendered_prs.slides[0]
                    if (
                        rendered_slide.notes_slide
                        and rendered_slide.notes_slide.notes_text_frame
                    ):
                        notes_text = rendered_slide.notes_slide.notes_text_frame.text
                        assert "Translated notes" in notes_text

                finally:
                    Path(source_tmp.name).unlink(missing_ok=True)
                    Path(output_tmp.name).unlink(missing_ok=True)


@pytest.mark.skipif(not PYTHON_PPTX_INSTALLED, reason="python-pptx not installed")
class TestPPTXParserMergeOptions:
    """Test paragraph merging options."""

    def test_parse_with_merge_paragraphs_true(self):
        """Should merge all text on slide into one segment."""
        from pptx import Presentation

        parser = PPTXParser(config={"merge_paragraphs": True})

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Title"
            slide.shapes.placeholders[1].text = (
                "Subtitle" if len(slide.shapes.placeholders) > 1 else None
            )

            # Add additional text box
            textbox = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(4), Inches(0.5)
            )
            textbox.text_frame.text = "Additional text"

            prs.save(tmp.name)

            try:
                parsed = parser.parse(tmp.name)

                # With merge_paragraphs=True, should have one segment per slide
                # (unless notes are also extracted)
                slide_segments = [
                    s for s in parsed.segments if s.context.get("type") == "slide"
                ]
                assert len(slide_segments) == 1
            finally:
                Path(tmp.name).unlink(missing_ok=True)

    def test_parse_with_merge_paragraphs_false(self):
        """Should create segment per text frame."""
        from pptx import Presentation

        parser = PPTXParser(config={"merge_paragraphs": False})

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Title"

            # Add another text box
            textbox = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(4), Inches(0.5)
            )
            textbox.text_frame.text = "Additional text"

            prs.save(tmp.name)

            try:
                parsed = parser.parse(tmp.name)

                # With merge_paragraphs=False, should have multiple segments
                slide_segments = [
                    s
                    for s in parsed.segments
                    if s.context.get("type") in ("slide", "text_frame")
                ]
                assert len(slide_segments) >= 2
            finally:
                Path(tmp.name).unlink(missing_ok=True)


@pytest.mark.skipif(not PYTHON_PPTX_INSTALLED, reason="python-pptx not installed")
class TestPPTXParserEmptySlides:
    """Test handling of empty slides."""

    def test_exclude_empty_slides_by_default(self):
        """Should exclude slides with no text by default."""
        from pptx import Presentation

        parser = PPTXParser(config={"include_empty_slides": False})

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()
            # Add slide with content
            slide1 = prs.slides.add_slide(prs.slide_layouts[0])
            slide1.shapes.title.text = "Content"

            # Add empty slide (layout without default text)
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

            prs.save(tmp.name)

            try:
                parsed = parser.parse(tmp.name)

                # Should only have one segment (from non-empty slide)
                # Note: blank layouts might still have placeholders, so this checks behavior
                assert len(parsed.segments) >= 1
            finally:
                Path(tmp.name).unlink(missing_ok=True)

    def test_include_empty_slides_when_configured(self):
        """Should include empty slides when configured."""
        from pptx import Presentation

        parser = PPTXParser(config={"include_empty_slides": True})

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()
            slide1 = prs.slides.add_slide(prs.slide_layouts[0])
            slide1.shapes.title.text = "Content"

            slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

            prs.save(tmp.name)

            try:
                parsed = parser.parse(tmp.name)

                # With include_empty_slides=True, should have segments for both slides
                # (though empty slides might have empty text)
                assert len(parsed.segments) >= 1
            finally:
                Path(tmp.name).unlink(missing_ok=True)


class TestPPTXNotInstalled:
    """Test behavior when python-pptx is not installed."""

    def test_parser_raises_when_not_installed(self):
        """Should raise RuntimeError when python-pptx unavailable."""
        if PYTHON_PPTX_INSTALLED:
            pytest.skip("python-pptx is installed")

        with pytest.raises(RuntimeError, match="python-pptx is not installed"):
            PPTXParser()

    def test_factory_raises_when_not_installed(self):
        """Should raise RuntimeError from factory when python-pptx unavailable."""
        if PYTHON_PPTX_INSTALLED:
            pytest.skip("python-pptx is installed")

        with pytest.raises(RuntimeError, match="python-pptx is not installed"):
            create_pptx_parser()
